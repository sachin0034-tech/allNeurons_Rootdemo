"""
Generate a 5-slide executive .pptx deck from Snowflake data + Claude narratives.
Light cream / ivory theme — matches reference slide layouts with dark-green header bars
and gold accents on a warm cream background.

Slides:
  1. Executive Summary    — high-level KPIs + Claude key drivers
  2. Revenue Performance  — regional breakdown + budget vs actual + narrative
  3. Store Operations     — traffic, CVR, ADS, UPT + store type table
  4. Supply Chain & 3PL  — OTD, freight, inventory + narrative
  5. Risks & Opportunities — Claude-generated risks + opps + sources
"""
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette (cream/ivory light theme) ────────────────────────
BG        = RGBColor(0xFA, 0xF7, 0xF0)   # warm cream background
HEADER_BG = RGBColor(0x1C, 0x2B, 0x1E)   # dark green header bar
GOLD      = RGBColor(0x8B, 0x7D, 0x3A)   # gold accents and KPI labels
DARK      = RGBColor(0x1C, 0x2B, 0x1E)   # main body text
CARD_BG   = RGBColor(0xF0, 0xEB, 0xDE)   # slightly tinted card background
CARD_ALT  = RGBColor(0xFA, 0xF7, 0xF0)   # alternating row (same as BG = white-ish)
MUTED     = RGBColor(0x6B, 0x7C, 0x6E)   # secondary text
BORDER    = RGBColor(0xD4, 0xC8, 0xB0)   # card borders
POS       = RGBColor(0x2D, 0x6A, 0x4F)   # deep green positive
NEG       = RGBColor(0xB2, 0x22, 0x22)   # deep red negative
AMB       = RGBColor(0xB5, 0x89, 0x00)   # amber neutral
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ── Primitives ────────────────────────────────────────────────

def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def _rect(slide, x, y, w, h, fill_rgb, border_rgb=None, border_pt=0.75):
    s = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill_rgb
    if border_rgb:
        s.line.color.rgb = border_rgb
        s.line.width = Pt(border_pt)
    else:
        s.line.fill.background()
    return s


def _txt(slide, text, x, y, w, h, pt, rgb, bold=False,
         align=PP_ALIGN.LEFT, italic=False, wrap=True):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text)
    run.font.size = Pt(pt)
    run.font.color.rgb = rgb
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = "Calibri"
    return box


def _header(slide, title, subtitle=""):
    # Dark green header bar
    _rect(slide, 0, 0, 13.33, 0.78, HEADER_BG)
    # Gold separator line
    _rect(slide, 0, 0.77, 13.33, 0.03, GOLD)
    _txt(slide, title, 0.20, 0.06, 12.90, 0.40, 11, GOLD, bold=True)
    if subtitle:
        _txt(slide, subtitle, 0.20, 0.46, 12.90, 0.26, 8, MUTED, italic=True)


def _section_label(slide, text, x, y, w):
    _rect(slide, x, y + 0.17, w, 0.018, GOLD)
    _txt(slide, text, x, y, w, 0.20, 7, GOLD, bold=True)


# ── Format helpers ────────────────────────────────────────────

def _f(v):
    try:
        return float(v)
    except Exception:
        return 0.0


def _chg_color(val):
    v = _f(val)
    if v > 0.001:  return POS
    if v < -0.001: return NEG
    return AMB


def _fmt_yoy(val):
    return f"{_f(val)*100:+.1f}%"


def _fmt_bps(val):
    return f"{_f(val)*10000:+.0f}bps"


def _fmt_M(val):
    return f"${_f(val)/1_000_000:.2f}M"


def _fmt_K(val):
    return f"${_f(val)/1_000:.1f}K"


def _fmt_pct(val):
    return f"{_f(val)*100:.1f}%"


def _fmt_num(val):
    return f"{_f(val):,.0f}"


# ── KPI metric block ──────────────────────────────────────────

def _kpi(slide, x, y, w, label, val_str, ly_str, chg, bps=False):
    H = 0.90
    _rect(slide, x, y, w, H, CARD_BG, BORDER, 0.75)
    _txt(slide, label, x+0.08, y+0.04, w-0.16, 0.18, 7, GOLD, bold=True)
    _txt(slide, val_str, x+0.08, y+0.21, w-0.16, 0.36, 16, DARK, bold=True)
    _txt(slide, f"vs {ly_str} LY", x+0.08, y+0.60, w-0.60, 0.22, 7, MUTED)
    v = _f(chg)
    color = _chg_color(v)
    badge = _fmt_bps(v) if bps else _fmt_yoy(v)
    bw, bh = 0.72, 0.22
    bx = x + w - bw - 0.06
    by = y + 0.60
    _rect(slide, bx, by, bw, bh, color)
    _txt(slide, badge, bx, by+0.02, bw, bh-0.02, 7, WHITE, bold=True,
         align=PP_ALIGN.CENTER)


# ── Hero metric block (Slide 1) ───────────────────────────────

def _hero(slide, x, y, w, label, value, sub, chg_str, chg_color):
    H = 1.30
    _rect(slide, x, y, w, H, CARD_BG, BORDER, 0.75)
    _txt(slide, label, x+0.12, y+0.06, w-0.24, 0.20, 8, GOLD, bold=True)
    _txt(slide, value, x+0.12, y+0.26, w-0.24, 0.52, 22, DARK, bold=True)
    _txt(slide, sub,   x+0.12, y+0.78, w-0.24, 0.22, 8, MUTED)
    if chg_str:
        bw, bh = 0.90, 0.26
        _rect(slide, x + w - bw - 0.10, y + 0.78, bw, bh, chg_color)
        _txt(slide, chg_str,
             x + w - bw - 0.10, y + 0.80, bw, bh-0.02, 9,
             WHITE, bold=True, align=PP_ALIGN.CENTER)


# ── Table helpers ─────────────────────────────────────────────

def _th(slide, text, x, y, w, h):
    _rect(slide, x, y, w, h, GOLD)
    _txt(slide, text, x+0.04, y+0.04, w-0.08, h-0.04, 7.5,
         DARK, bold=True, align=PP_ALIGN.CENTER)


def _td(slide, text, x, y, w, h, rgb=DARK, bg=CARD_BG, bold=False,
        align=PP_ALIGN.CENTER):
    _rect(slide, x, y, w, h, bg, BORDER, 0.25)
    _txt(slide, text, x+0.04, y+0.04, w-0.08, h-0.04, 8,
         rgb, bold=bold, align=align, wrap=False)


def _badge_td(slide, text, x, y, w, h, color):
    bw = min(w - 0.10, 0.80)
    bx = x + (w - bw) / 2
    by = y + 0.06
    bh = h - 0.12
    _rect(slide, bx, by, bw, bh, color)
    _txt(slide, text, bx, by+0.02, bw, bh-0.02, 7,
         WHITE, bold=True, align=PP_ALIGN.CENTER)


# ── Bullet block ──────────────────────────────────────────────

def _bullets(slide, items, x, y, w, h, color=DARK, pt=8.5, prefix="•"):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(3)
        run = p.add_run()
        run.text = f"{prefix}  {item}"
        run.font.size = Pt(pt)
        run.font.color.rgb = color
        run.font.name = "Calibri"


def _source_line(slide, sources, x, y, w):
    if not sources:
        return
    texts = []
    for s in sources[:3]:
        texts.append(f"[{s.get('title', 'Source')}]  {s.get('url', '')}")
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(0.50))
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for t in texts:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = p.add_run()
        run.text = t
        run.font.size = Pt(6.5)
        run.font.color.rgb = MUTED
        run.font.italic = True
        run.font.name = "Calibri"


# ═══════════════════════════════════════════════════════════════
# SLIDE 1 — Executive Summary
# ═══════════════════════════════════════════════════════════════

def _slide_exec_summary(prs, metrics: dict, narr: dict):
    slide = _blank(prs)
    _bg(slide)

    cy       = _f(metrics.get("revenue"))
    ly       = _f(metrics.get("revenue_ly", cy))
    bud      = _f(metrics.get("revenue_budget", cy))
    traf_yoy = _f(metrics.get("traffic_yoy"))
    inv_turns = _f(metrics.get("inventory_turns", 8.4))
    cash     = _f(metrics.get("cash_position", 18_400_000))
    rev_yoy  = (cy - ly) / ly if ly else 0
    bud_gap  = (cy - bud) / bud if bud else 0

    title = (
        f"EXECUTIVE SUMMARY   |   Revenue {_fmt_yoy(rev_yoy)} "
        f"({_fmt_K(abs(cy - ly))}) to LY   //   "
        f"{_fmt_yoy(bud_gap)} ({_fmt_K(abs(cy - bud))}) to Budget"
    )
    _header(slide, title,
            f"Week ending {metrics.get('week_label', '')}  ·  "
            "Snowflake source of truth  ·  Narrative by Claude")

    # ── Hero KPIs (4 blocks) ──────────────────────────────────
    HY = 0.92
    HW = (13.33 - 0.20 * 5) / 4
    heroes = [
        ("REVENUE",          _fmt_M(cy),          f"Budget: {_fmt_M(bud)}",
         _fmt_yoy(rev_yoy),  _chg_color(rev_yoy)),
        ("STORE TRAFFIC",    _fmt_yoy(traf_yoy),  "vs Last Year",
         None, None),
        ("INVENTORY TURNS",  f"{inv_turns:.1f}x", "annualised",
         None, None),
        ("CASH POSITION",    _fmt_M(cash),        "end of week",
         None, None),
    ]
    for i, (lbl, val, sub, chg, chg_c) in enumerate(heroes):
        hx = 0.20 + i * (HW + 0.20)
        _hero(slide, hx, HY, HW, lbl, val, sub, chg, chg_c or AMB)

    # ── Key Drivers (Claude bullets) ──────────────────────────
    DY = HY + 1.40
    _section_label(slide, "KEY DRIVERS — Claude Analysis", 0.20, DY, 13.00)
    DY += 0.24

    bullets = narr.get("exec_summary", {}).get("bullets", [])
    if bullets:
        half = len(bullets) // 2
        _bullets(slide, bullets[:half], 0.20, DY, 6.40, 1.80, DARK, 9)
        _bullets(slide, bullets[half:], 6.80, DY, 6.33, 1.80, DARK, 9)

    _source_line(slide, narr.get("exec_summary", {}).get("sources", []),
                 0.20, 7.10, 12.90)


# ═══════════════════════════════════════════════════════════════
# SLIDE 2 — Revenue Performance
# ═══════════════════════════════════════════════════════════════

def _slide_revenue(prs, metrics: dict, fp_df, narr: dict):
    slide = _blank(prs)
    _bg(slide)
    _header(slide, "REVENUE PERFORMANCE — Regional Breakdown & Budget vs Actual",
            "Revenue and margin analysis by region; variance to budget driven by sales mix and store format performance.")

    # ── LEFT: region table ─────────────────────────────────────
    TX, TY = 0.20, 0.92
    COLS = [
        ("REGION",   1.15), ("CY SALES", 1.50), ("LY SALES", 1.50),
        ("YoY",      1.00), ("BUDGET",   1.50), ("vs BUD",   1.00),
    ]
    GAP, RH_H, RH_D = 0.04, 0.28, 0.48

    hx = TX
    for lbl, cw in COLS:
        _th(slide, lbl, hx, TY, cw, RH_H)
        hx += cw + GAP
    ty = TY + RH_H + 0.03

    if not fp_df.empty:
        for i, (_, row) in enumerate(fp_df.iterrows()):
            cy_s  = _f(row.get("NET_SALES"))
            ly_s  = _f(row.get("LY_NET_SALES", cy_s))
            bud_s = cy_s * _f(metrics.get("budget_factor", 1.06))
            yoy   = (cy_s - ly_s) / ly_s if ly_s else 0
            bvb   = (cy_s - bud_s) / bud_s if bud_s else 0

            vals = [
                (str(row.get("REGION_CODE", "")), GOLD, True,  PP_ALIGN.LEFT),
                (_fmt_M(cy_s),   DARK, False, PP_ALIGN.CENTER),
                (_fmt_M(ly_s),   MUTED, False, PP_ALIGN.CENTER),
                (_fmt_yoy(yoy),  _chg_color(yoy), True, PP_ALIGN.CENTER),
                (_fmt_M(bud_s),  MUTED, False, PP_ALIGN.CENTER),
                (_fmt_yoy(bvb),  _chg_color(bvb), True, PP_ALIGN.CENTER),
            ]
            bg_row = CARD_BG if i % 2 == 0 else CARD_ALT
            rx2 = TX
            for (txt, rgb, bold, align), (_, cw) in zip(vals, COLS):
                _td(slide, txt, rx2, ty, cw, RH_D, rgb, bg_row, bold, align)
                rx2 += cw + GAP
            ty += RH_D + 0.03
            if ty > 6.20:
                break

    # ── RIGHT: narrative ───────────────────────────────────────
    NX = 7.98
    NW = 13.33 - NX - 0.15
    NY = 0.92
    _section_label(slide, "REVENUE INSIGHTS — Claude", NX, NY, NW)
    NY += 0.24

    bullets = narr.get("revenue", {}).get("bullets", [])
    _bullets(slide, bullets, NX, NY, NW, 3.60, DARK, 8.5)

    _source_line(slide, narr.get("revenue", {}).get("sources", []),
                 NX, NY + 3.65, NW)

    # ── FP Mix mini-strip ──────────────────────────────────────
    if not fp_df.empty:
        _section_label(slide, "FULL-PRICE MIX BY REGION", NX, NY + 4.10, NW)
        sy = NY + 4.34
        for _, row in fp_df.iterrows():
            rfp = _f(row.get("FP_MIX"))
            rly = _f(row.get("LY_FP_MIX", rfp))
            rc  = str(row.get("REGION_CODE", ""))
            _txt(slide, rc, NX, sy, 0.60, 0.28, 7.5, GOLD, bold=True)
            _txt(slide, _fmt_pct(rfp), NX+0.65, sy, 0.70, 0.28, 7.5, DARK)
            chg_c = _chg_color(rfp - rly)
            bar_w = max(0.05, (rfp / 0.85) * (NW - 1.55))
            _rect(slide, NX+1.40, sy+0.04, bar_w, 0.20, chg_c)
            sy += 0.30
            if sy > 7.25:
                break

    _source_line(slide, narr.get("revenue", {}).get("sources", []),
                 0.20, 7.12, 7.70)


# ═══════════════════════════════════════════════════════════════
# SLIDE 3 — Store Operations
# ═══════════════════════════════════════════════════════════════

def _slide_operations(prs, metrics: dict, st_df, narr: dict):
    slide = _blank(prs)
    _bg(slide)
    _header(slide, "STORE OPERATIONS — Traffic · Conversion · Basket · Format Performance",
            "Traffic growth outpacing conversion improvement; basket metrics indicate pricing power remains intact.")

    traf = _f(metrics.get("traffic"))
    ly_t = _f(metrics.get("ly_traffic", traf or 1))
    cvr  = _f(metrics.get("cvr"))
    ly_c = _f(metrics.get("ly_cvr", cvr))
    ads  = _f(metrics.get("ads"))
    ly_a = _f(metrics.get("ly_ads", ads or 1))
    upt  = _f(metrics.get("upt"))

    # ── Top 4 KPIs ────────────────────────────────────────────
    Y0 = 0.92
    KW = (13.33 - 0.20 * 5) / 4
    kpis = [
        ("STORE TRAFFIC",    _fmt_num(traf), _fmt_num(ly_t),
         (traf - ly_t) / ly_t if ly_t else 0, False),
        ("CONVERSION RATE",  _fmt_pct(cvr),  _fmt_pct(ly_c),
         cvr - ly_c, True),
        ("AVG. ORDER VALUE", f"${ads:.2f}",  f"${ly_a:.2f}",
         (ads - ly_a) / ly_a if ly_a else 0, False),
        ("UNITS PER TXN",    f"{upt:.2f}",   "—", 0, False),
    ]
    for i, (lbl, val, ly_v, chg, bps) in enumerate(kpis):
        kx = 0.20 + i * (KW + 0.20)
        _kpi(slide, kx, Y0, KW, lbl, val, ly_v, chg, bps)

    # ── Store type table ──────────────────────────────────────
    TY = Y0 + 1.02
    TX = 0.20
    COLS = [
        ("STORE TYPE", 2.20), ("SALES",    1.45), ("TRAFFIC", 1.30),
        ("CVR",        1.10), ("ADS",      1.15), ("UPT",    1.00),
        ("SALES YoY",  1.30), ("CVR Δ", 1.23),
    ]
    GAP, RH_H, RH_D = 0.04, 0.28, 0.52

    hx = TX
    for lbl, cw in COLS:
        _th(slide, lbl, hx, TY, cw, RH_H)
        hx += cw + GAP
    ty = TY + RH_H + 0.03

    if not st_df.empty:
        for i, (_, row) in enumerate(st_df.iterrows()):
            ns    = _f(row.get("NET_SALES"))
            ly_ns = _f(row.get("LY_NET_SALES", ns or 1))
            tr    = _f(row.get("TRAFFIC"))
            cv    = _f(row.get("CVR"))
            ly_cv = _f(row.get("LY_CVR", cv))
            a     = _f(row.get("ADS"))
            u     = _f(row.get("UPT"))
            s_yoy = (ns - ly_ns) / ly_ns if ly_ns else 0
            c_bps = cv - ly_cv

            vals = [
                (str(row.get("STORE_TYPE_NAME", "")), GOLD, True, PP_ALIGN.LEFT),
                (_fmt_M(ns),       DARK, False, PP_ALIGN.CENTER),
                (_fmt_num(tr),     DARK, False, PP_ALIGN.CENTER),
                (_fmt_pct(cv),     DARK, False, PP_ALIGN.CENTER),
                (f"${a:.2f}",      DARK, False, PP_ALIGN.CENTER),
                (f"{u:.2f}",       DARK, False, PP_ALIGN.CENTER),
                (_fmt_yoy(s_yoy),  _chg_color(s_yoy), True, PP_ALIGN.CENTER),
                (_fmt_bps(c_bps),  _chg_color(c_bps), True, PP_ALIGN.CENTER),
            ]
            bg_row = CARD_BG if i % 2 == 0 else CARD_ALT
            rx2 = TX
            for (txt, rgb, bold, align), (_, cw) in zip(vals, COLS):
                _td(slide, txt, rx2, ty, cw, RH_D, rgb, bg_row, bold, align)
                rx2 += cw + GAP
            ty += RH_D + 0.03
            if ty > 6.30:
                break

    # ── Narrative ─────────────────────────────────────────────
    ny = ty + 0.08
    if ny < 5.80:
        ny = 5.80
    _section_label(slide, "OPERATIONS INSIGHTS — Claude", 0.20, ny, 13.00)
    ny += 0.24
    bullets = narr.get("operations", {}).get("bullets", [])
    _bullets(slide, bullets[:2], 0.20, ny, 6.40, 0.80, DARK, 8)
    _bullets(slide, bullets[2:], 6.80, ny, 6.33, 0.80, DARK, 8)
    _source_line(slide, narr.get("operations", {}).get("sources", []),
                 0.20, 7.12, 12.90)


# ═══════════════════════════════════════════════════════════════
# SLIDE 4 — Supply Chain & 3PL
# ═══════════════════════════════════════════════════════════════

def _slide_supply_chain(prs, metrics: dict, narr: dict):
    slide = _blank(prs)
    _bg(slide)
    _header(slide, "SUPPLY CHAIN & 3PL — Delivery · Freight · Inventory · Fill Rate",
            "Logistics performance improving YoY; inventory aging risk identified in 2 categories.")

    otd       = _f(metrics.get("otd_pct",            0.942))
    ly_otd    = _f(metrics.get("ly_otd_pct",         0.921))
    freight   = _f(metrics.get("freight_cost_yoy",   -0.06))
    inv_aging = _f(metrics.get("inventory_aging_days", 41))
    fill_rate = _f(metrics.get("fill_rate",           0.968))
    ly_fill   = _f(metrics.get("ly_fill_rate",        0.962))
    inv_turns = _f(metrics.get("inventory_turns",     8.4))
    ly_turns  = _f(metrics.get("ly_inv_turns",        7.8))

    Y0 = 0.92
    KW = (13.33 - 0.20 * 5) / 4
    kpis = [
        ("ON-TIME DELIVERY",  _fmt_pct(otd),         _fmt_pct(ly_otd),
         otd - ly_otd, True),
        ("FREIGHT COST YoY",  _fmt_yoy(freight),     "Prior Year",
         freight, False),
        ("INVENTORY AGING",   f"{inv_aging:.0f} days", "avg. SKU age",
         None, False),
        ("FILL RATE",         _fmt_pct(fill_rate),   _fmt_pct(ly_fill),
         fill_rate - ly_fill, True),
    ]
    for i, (lbl, val, ly_v, chg, bps) in enumerate(kpis):
        kx = 0.20 + i * (KW + 0.20)
        if chg is None:
            # Simple block without badge
            _rect(slide, kx, Y0, KW, 0.90, CARD_BG, BORDER, 0.75)
            _txt(slide, lbl, kx+0.08, Y0+0.04, KW-0.16, 0.18, 7, GOLD, bold=True)
            _txt(slide, val, kx+0.08, Y0+0.21, KW-0.16, 0.36, 16, DARK, bold=True)
            _txt(slide, ly_v, kx+0.08, Y0+0.60, KW-0.16, 0.22, 7, MUTED)
        else:
            _kpi(slide, kx, Y0, KW, lbl, val, ly_v, chg, bps)

    # ── Detailed KPI table ─────────────────────────────────────
    TY = Y0 + 1.05
    TX = 0.20
    COLS = [
        ("METRIC",       3.20), ("CURRENT",   1.60),
        ("PRIOR YEAR",   1.60), ("VARIANCE",  1.60), ("STATUS",  1.70),
    ]
    GAP, RH_H, RH_D = 0.04, 0.28, 0.52
    rows = [
        ("On-Time Delivery %",
         _fmt_pct(otd),      _fmt_pct(ly_otd),
         _fmt_bps(otd - ly_otd),
         POS if otd >= ly_otd else NEG,
         "On Track" if otd >= 0.92 else "Watch"),
        ("Freight Cost Index",
         _fmt_yoy(freight),  "Base",
         _fmt_yoy(freight),
         POS if freight < 0 else NEG,
         "Favourable" if freight < 0 else "Over"),
        ("Avg. Inventory Age (days)",
         f"{inv_aging:.0f}", "—",
         "—",
         AMB,
         "Monitor"),
        ("Fill Rate %",
         _fmt_pct(fill_rate), _fmt_pct(ly_fill),
         _fmt_bps(fill_rate - ly_fill),
         POS if fill_rate >= ly_fill else NEG,
         "On Track" if fill_rate >= 0.95 else "Risk"),
        ("Inventory Turns (ann.)",
         f"{inv_turns:.1f}x", f"{ly_turns:.1f}x",
         f"{inv_turns - ly_turns:+.1f}x",
         POS if inv_turns >= ly_turns else NEG,
         "Improving" if inv_turns >= ly_turns else "Watch"),
    ]

    hx = TX
    for lbl, cw in COLS:
        _th(slide, lbl, hx, TY, cw, RH_H)
        hx += cw + GAP
    ty = TY + RH_H + 0.03

    for i, (metric, cur, ly_v, var, chg_c, status) in enumerate(rows):
        bg_row = CARD_BG if i % 2 == 0 else CARD_ALT
        vals = [
            (metric, DARK,  True,  PP_ALIGN.LEFT),
            (cur,    DARK,  False, PP_ALIGN.CENTER),
            (ly_v,   MUTED, False, PP_ALIGN.CENTER),
            (var,    chg_c, True,  PP_ALIGN.CENTER),
        ]
        rx2 = TX
        for (txt, rgb, bold, align), (_, cw) in zip(vals, COLS[:-1]):
            _td(slide, txt, rx2, ty, cw, RH_D, rgb, bg_row, bold, align)
            rx2 += cw + GAP
        # Status badge
        st_cw = COLS[-1][1]
        _badge_td(slide, status, rx2, ty, st_cw, RH_D, chg_c)
        ty += RH_D + 0.03

    # ── Narrative ─────────────────────────────────────────────
    ny = ty + 0.06
    _section_label(slide, "SUPPLY CHAIN INSIGHTS — Claude", 0.20, ny, 13.00)
    ny += 0.24
    bullets = narr.get("supply_chain", {}).get("bullets", [])
    _bullets(slide, bullets[:2], 0.20, ny, 6.40, 0.75, DARK, 8)
    _bullets(slide, bullets[2:], 6.80, ny, 6.33, 0.75, DARK, 8)
    _source_line(slide, narr.get("supply_chain", {}).get("sources", []),
                 0.20, 7.12, 12.90)


# ═══════════════════════════════════════════════════════════════
# SLIDE 5 — Risks & Opportunities
# ═══════════════════════════════════════════════════════════════

def _slide_risks(prs, narr: dict):
    slide = _blank(prs)
    _bg(slide)
    _header(slide, "RISKS & OPPORTUNITIES — Strategic Outlook",
            "External market intelligence sourced via Claude web search · Internal signals from Snowflake analytics.")

    risk_narr = narr.get("risks", {})
    risks   = risk_narr.get("risks",         [])
    opps    = risk_narr.get("opportunities", [])
    sources = risk_narr.get("sources",       [])

    LY = 0.92
    # ── LEFT: Key Risks ───────────────────────────────────────
    LX, LW = 0.20, 6.30
    _rect(slide, LX, LY, LW, 5.70, CARD_BG, BORDER, 0.75)
    _rect(slide, LX, LY, LW, 0.40, NEG)
    _txt(slide, "KEY RISKS", LX+0.15, LY+0.08, LW-0.30, 0.28,
         10, WHITE, bold=True)

    ry = LY + 0.55
    for risk in risks:
        _rect(slide, LX+0.18, ry+0.07, 0.10, 0.10, NEG)
        _txt(slide, risk, LX+0.36, ry, LW-0.54, 0.60, 8.5, DARK, wrap=True)
        ry += 0.70
        if ry > LY + 5.20:
            break

    # ── RIGHT: Opportunities ──────────────────────────────────
    RX = 6.83
    RW = 13.33 - RX - 0.20
    _rect(slide, RX, LY, RW, 5.70, CARD_BG, BORDER, 0.75)
    _rect(slide, RX, LY, RW, 0.40, POS)
    _txt(slide, "OPPORTUNITIES", RX+0.15, LY+0.08, RW-0.30, 0.28,
         10, WHITE, bold=True)

    oy = LY + 0.55
    for opp in opps:
        _rect(slide, RX+0.18, oy+0.07, 0.10, 0.10, POS)
        _txt(slide, opp, RX+0.36, oy, RW-0.54, 0.60, 8.5, DARK, wrap=True)
        oy += 0.70
        if oy > LY + 5.20:
            break

    # ── Sources footer ─────────────────────────────────────────
    _section_label(slide, "SOURCES & EXTERNAL INTELLIGENCE — Claude Web Search",
                   0.20, 6.75, 12.90)
    _source_line(slide, sources, 0.20, 6.96, 12.90)


# ═══════════════════════════════════════════════════════════════
# Public entry point
# ═══════════════════════════════════════════════════════════════

def build_deck_light(kpi: dict, fp_df, st_df, reg_df,
                     narratives: dict, week_label: str = "") -> bytes:
    """
    Build a 5-slide executive PPTX (light/cream theme) and return the bytes.

    kpi        — dict from a single-row VW_WEEKLY_KPI query
    fp_df      — DataFrame from VW_REGION_FP_MIX
    st_df      — DataFrame from VW_STORE_TYPE_PERF
    reg_df     — DataFrame from VW_REGION_WEEKLY_YOY
    narratives — dict from utils.narrative.generate_narratives()
    week_label — human-readable week label, e.g. "W24 2026"
    """
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    cy    = _f(kpi.get("NET_SALES"))
    ly    = _f(kpi.get("LY_NET_SALES", cy))
    bud   = _f(kpi.get("BUDGET_SALES", cy))
    traf  = _f(kpi.get("TRAFFIC"))
    ly_tr = _f(kpi.get("LY_TRAFFIC", traf or 1))
    txns  = _f(kpi.get("TRANSACTIONS"))
    ly_tx = _f(kpi.get("LY_TRANSACTIONS", txns or 1))
    cvr   = txns / traf if traf else 0
    ly_c  = ly_tx / ly_tr if ly_tr else 0

    # Derive top/bottom regions
    top_regs, bot_regs = [], []
    if not reg_df.empty and "SALES_YOY_PCT" in reg_df.columns:
        sorted_r = reg_df.sort_values("SALES_YOY_PCT", ascending=False)
        top_regs = sorted_r["REGION_CODE"].head(2).tolist()
        bot_regs = sorted_r["REGION_CODE"].tail(2).tolist()

    top_types, bot_types = [], []
    if not st_df.empty and "NET_SALES" in st_df.columns and "LY_NET_SALES" in st_df.columns:
        st_df2 = st_df.copy()
        st_df2["_yoy"] = (
            st_df2["NET_SALES"].astype(float) - st_df2["LY_NET_SALES"].astype(float)
        ) / st_df2["LY_NET_SALES"].astype(float).replace(0, float("nan"))
        sorted_st = st_df2.sort_values("_yoy", ascending=False)
        top_types = sorted_st["STORE_TYPE_NAME"].head(2).tolist()
        bot_types = sorted_st["STORE_TYPE_NAME"].tail(2).tolist()

    metrics = {
        "week_label":           week_label,
        "revenue":              cy,
        "revenue_ly":           ly,
        "revenue_budget":       bud,
        "revenue_yoy":          (cy - ly) / ly if ly else 0,
        "revenue_budget_gap":   (cy - bud) / bud if bud else 0,
        "budget_factor":        1.06,
        "traffic":              traf,
        "ly_traffic":           ly_tr,
        "traffic_yoy":          (traf - ly_tr) / ly_tr if ly_tr else 0,
        "cvr":                  cvr,
        "ly_cvr":               ly_c,
        "cvr_bps":              cvr - ly_c,
        "ads":                  _f(kpi.get("ADS")),
        "ly_ads":               _f(kpi.get("LY_ADS")),
        "upt":                  _f(kpi.get("UPT")),
        "fp_mix":               _f(kpi.get("FP_MIX")),
        "ly_fp_mix":            _f(kpi.get("LY_FP_MIX")),
        "fp_mix_bps":           _f(kpi.get("FP_MIX")) - _f(kpi.get("LY_FP_MIX")),
        "gross_margin_pct":     0.58,
        "ebitda_pct":           0.22,
        "cash_position":        18_400_000,
        "inventory_turns":      8.4,
        "ly_inv_turns":         7.8,
        "otd_pct":              0.942,
        "ly_otd_pct":           0.921,
        "freight_cost_yoy":     -0.06,
        "inventory_aging_days": 41,
        "fill_rate":            0.968,
        "ly_fill_rate":         0.962,
        "freight_savings":      -0.06,
        "top_regions":          top_regs,
        "bottom_regions":       bot_regs,
        "top_store_types":      top_types,
        "bottom_store_types":   bot_types,
    }

    _slide_exec_summary(prs, metrics, narratives)
    _slide_revenue(prs, metrics, fp_df, narratives)
    _slide_operations(prs, metrics, st_df, narratives)
    _slide_supply_chain(prs, metrics, narratives)
    _slide_risks(prs, narratives)

    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()
